# Build a real .pptx from Manus's 20 rendered slide images.
# Each slide = one 16:9 page with the full-bleed rendered image.
import glob
import os
import tempfile
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

SLIDES = r'D:\study\AI\manus-out\slides'
OUT = r'D:\study\AI\manus-out\未来五年（2025-2030）半导体行业发展趋势.pptx'

files = sorted(glob.glob(os.path.join(SLIDES, '*.webp')))
assert len(files) == 20, f'expected 20 slides, got {len(files)}'

prs = Presentation()
# 16:9: 13.333in x 7.5in
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]  # blank layout

# python-pptx does not accept webp directly -> convert to PNG first.
tmp = tempfile.mkdtemp(prefix='dsh-manus-pptx-')
png_files = []
for f in files:
    png = os.path.join(tmp, os.path.splitext(os.path.basename(f))[0] + '.png')
    Image.open(f).convert('RGB').save(png, 'PNG')
    png_files.append(png)

for png in png_files:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)

prs.save(OUT)
print('saved:', OUT)
print('slides:', len(prs.slides._sldIdLst))

